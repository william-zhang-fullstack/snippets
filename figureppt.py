import io
import pptx


class FigurePptSink:
    def __init__(self, height: float = 7.5, width: float = 10.) -> None:
        """ High level interface to write Matplotlib figs to PPT

        :param height: presentation height in inches, defaults to 7.5
        :type height: float, optional
        :param width: presentation width in inches, defaults to 10
        :type width: float, optional
        """
        self.buffer = io.BytesIO()
        self.height = height
        self.width = width
        ppt = pptx.Presentation()
        ppt.slide_height = pptx.util.Inches(height)
        ppt.slide_width = pptx.util.Inches(width)
        self.ppt = ppt

    def add_slide(self, fig) -> None:
        """ Append figure to running slide deck

        :param fig: figure to add
        :type  fig: matplotlib.pyplot.Figure
        """
        fig.savefig(self.buffer)  # load image into buffer
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[6])  # +blank
        slide.shapes.add_picture(
            self.buffer,
            left=pptx.util.Inches(0),
            top=pptx.util.Inches(0))  # add image to slide from buffer
        self.buffer.seek(0)  # move buff pointer to buff start
        self.buffer.truncate()  # clear image from buffer

    def save(self, filepath: str) -> None:
        """ Save running slide deck to disk

        :param filepath: target path on disk
        :type  filepath: str
        """
        self.ppt.save(filepath)
